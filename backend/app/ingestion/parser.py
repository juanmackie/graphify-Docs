"""Document text extraction for supported formats: PDF, DOCX, PPTX, TXT, MD, HTML."""
from __future__ import annotations

import re
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path

from pypdf import PdfReader

SUPPORTED_EXTS = {".pdf", ".docx", ".txt", ".md", ".pptx", ".html", ".htm"}


@dataclass
class ParsedDocument:
    text: str
    pages: int = 1
    format: str = ""

    @property
    def words(self) -> int:
        return len(self.text.split())


# ── HTML → plain text ─────────────────────────────────────────────────
_BLOCK_TAGS = {
    "p", "div", "h1", "h2", "h3", "h4", "h5", "h6", "li", "tr", "br",
    "section", "article", "blockquote", "pre", "ul", "ol", "table", "hr",
}


class _HTMLToText(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"script", "style"}:
            self._skip_depth += 1
        if tag in _BLOCK_TAGS and self.parts and not self.parts[-1].endswith("\n"):
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style"} and self._skip_depth:
            self._skip_depth -= 1
        if tag in _BLOCK_TAGS and self.parts and not self.parts[-1].endswith("\n"):
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._skip_depth:
            self.parts.append(data)


def html_to_text(raw: str) -> str:
    parser = _HTMLToText()
    parser.feed(raw)
    text = "".join(parser.parts)
    lines = [ln.strip() for ln in text.splitlines()]
    return "\n".join(ln for ln in lines if ln)


# ── per-format readers ────────────────────────────────────────────────
def _read_text_file(path: Path) -> str:
    """Read a plain-text file, trying common encodings."""
    raw = path.read_bytes()
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("latin-1", errors="replace")


def parse_pdf(path: Path) -> ParsedDocument:
    try:
        reader = PdfReader(str(path))
    except Exception as exc:  # noqa: BLE001
        raise ValueError(f"Could not open PDF: {exc}") from exc
    if reader.is_encrypted:
        raise ValueError("Encrypted PDFs are not supported — remove the password and retry.")
    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n\n".join(p.strip() for p in pages if p.strip())
    return ParsedDocument(text=text, pages=len(pages), format="pdf")


def parse_docx(path: Path) -> ParsedDocument:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover
        raise ValueError("python-docx is not installed.") from exc
    doc = Document(str(path))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts)
    if not text.strip():
        raise ValueError("No text found in the DOCX file.")
    return ParsedDocument(text=text, format="docx")


def parse_pptx(path: Path) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise ValueError("python-pptx is not installed.") from exc
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
    text = "\n".join(parts)
    if not text.strip():
        raise ValueError("No text found in the PPTX file.")
    return ParsedDocument(text=text, format="pptx")


def parse_txt(path: Path) -> ParsedDocument:
    text = _read_text_file(path).strip()
    if not text:
        raise ValueError("The text file is empty.")
    return ParsedDocument(text=text, format="txt")


def parse_html(path: Path) -> ParsedDocument:
    text = html_to_text(_read_text_file(path)).strip()
    if not text:
        raise ValueError("No readable text found in the HTML file.")
    return ParsedDocument(text=text, format="html")


_PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".txt": parse_txt,
    ".md": parse_txt,
    ".html": parse_html,
    ".htm": parse_html,
}


def parse_document(path: Path) -> ParsedDocument:
    """Extract plain text from *path* based on its file extension."""
    path = Path(path)
    ext = path.suffix.lower()
    if ext == ".doc":
        raise ValueError(
            "Legacy .doc files aren't supported — save the file as .docx (or .txt) and retry."
        )
    if ext not in SUPPORTED_EXTS:
        raise ValueError(f"Unsupported file type '{ext or '?'}'. Supported: {', '.join(sorted(SUPPORTED_EXTS))}")
    parsed = _PARSERS[ext](path)
    cleaned = re.sub(r"[ \t]+\n", "\n", parsed.text)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    parsed.text = cleaned.strip()
    if not parsed.text:
        raise ValueError("No extractable text found in the document.")
    return parsed
