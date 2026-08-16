"""Shared test fixtures: sample documents in every supported format."""
from __future__ import annotations

import os
import sys
import tempfile
from pathlib import Path

import pytest

BACKEND_ROOT = Path(__file__).resolve().parents[1]
if str(BACKEND_ROOT) not in sys.path:
    sys.path.insert(0, str(BACKEND_ROOT))

# Isolate app data (sqlite + uploads) from the repo — must be set before any
# `app.*` import (settings are read at import time).
TEST_DATA_DIR = Path(tempfile.mkdtemp(prefix="docgraph-test-"))
os.environ["DATA_DIR"] = str(TEST_DATA_DIR)

# Keep tests offline & deterministic: real env vars win over any .env file,
# so pin the LLM config to inert values (a developer's local .env must never
# trigger real API calls from the test suite).
os.environ["OPENAI_API_KEY"] = ""
os.environ["OPENAI_BASE_URL"] = "http://localhost:9/v1"
os.environ["OPENAI_MODEL"] = "test-model"

SAMPLE_TEXT = """Knowledge Graphs: An Overview

A knowledge graph represents a network of real-world entities and illustrates the
relationships between them. Search engines, social networks, and recommendation
systems all rely on knowledge graphs to power their features.

Entities are the fundamental units of a knowledge graph. Each entity has a type,
such as person, organization, or concept, and a set of attributes. Relationships
connect entities and carry a label that describes the connection.

Graph databases such as Neo4j store knowledge graphs efficiently and support fast
traversal queries. The property graph model used by Neo4j consists of nodes,
relationships, and properties.

Community detection groups related entities into clusters. The Leiden algorithm
improves on the Louvain method by guaranteeing well-connected communities.
Detecting communities helps analysts understand the overall structure of a graph.

Export formats let users share a knowledge graph with collaborators. Common
formats include JSON, GraphML, and CSV. An interactive HTML export is the most
accessible way to share a graph with non-technical stakeholders.

Visualization is critical. Force-directed layouts position connected nodes close
together, making clusters and bridges easy to spot. Interactive tooling lets
users search for an entity, highlight its neighbors, and trace paths between two
nodes.

This document itself is an example: it can be parsed, chunked, and mapped into a
knowledge graph showing how concepts like entities, relationships, communities,
and visualization connect to one another.
"""


def _make_pdf(path: Path, text: str) -> Path:
    from pypdf import PdfWriter
    from pypdf.generic import DictionaryObject, NameObject, StreamObject

    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    resources = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})}
    )
    page[NameObject("/Resources")] = resources
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    content = f"BT /F1 11 Tf 72 720 Td ({escaped[:2000]}) Tj ET".encode("latin-1")
    stream = StreamObject()
    stream.set_data(content)
    page[NameObject("/Contents")] = writer._add_object(stream)
    writer.write(str(path))
    return path


def _make_docx(path: Path, text: str) -> Path:
    from docx import Document

    doc = Document()
    for para in text.split("\n\n"):
        doc.add_paragraph(para)
    doc.save(str(path))
    return path


def _make_pptx(path: Path, text: str) -> Path:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = text.splitlines()[0]
    body = slide.placeholders[1].text_frame
    for line in text.splitlines()[1:]:
        body.add_paragraph().text = line
    prs.save(str(path))
    return path


@pytest.fixture
def sample_text() -> str:
    return SAMPLE_TEXT


@pytest.fixture
def sample_txt(tmp_path: Path) -> Path:
    path = tmp_path / "sample.txt"
    path.write_text(SAMPLE_TEXT, encoding="utf-8")
    return path


@pytest.fixture
def sample_md(tmp_path: Path) -> Path:
    path = tmp_path / "sample.md"
    path.write_text(f"# Knowledge Graphs\n\n{SAMPLE_TEXT}", encoding="utf-8")
    return path


@pytest.fixture
def sample_html(tmp_path: Path) -> Path:
    path = tmp_path / "sample.html"
    path.write_text(
        f"<html><head><title>KG</title><style>p{{color:red}}</style></head>"
        f"<body><h1>Knowledge Graphs</h1><p>{SAMPLE_TEXT.replace(chr(10), '</p><p>')}</p>"
        f"<script>alert('x')</script></body></html>",
        encoding="utf-8",
    )
    return path


@pytest.fixture
def sample_pdf(tmp_path: Path) -> Path:
    return _make_pdf(tmp_path / "sample.pdf", SAMPLE_TEXT)


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    return _make_docx(tmp_path / "sample.docx", SAMPLE_TEXT)


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    return _make_pptx(tmp_path / "sample.pptx", SAMPLE_TEXT)
